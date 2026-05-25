"""
Editorial 16-slide PowerPoint deck — matches the Spondon Board Presentation preview.

Slide map (16 slides, 16:9):
   1. Title cover (dark navy → blue gradient, italic "Programme update.")
   2. Agenda (numbered list of 5 sections)
   3. Section 01 divider — Key indicators
   4. Big number (huge italic total, +%MoM)
   5. Hero quote / context
   6. KPI dashboard (4-up cards)
   7. Section 02 divider — Activity by category
   8. Stacked category breakdown
   9. Top forms / activity ranking
  10. Section 03 divider — Geography & centres
  11. Top districts leaderboard
  12. Service centre breakdown
  13. Section 04 divider — Watch list & alerts
  14. Closing quote (coral gradient)
  15. Forward look / next month
  16. Thank you / contact (dark)
"""
from __future__ import annotations

import io
from datetime import date as _date

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Editorial palette ─────────────────────────────────────────────────────────
NAVY        = RGBColor(0x00, 0x2A, 0x3D)
UNFPA       = RGBColor(0x00, 0x65, 0x8C)
UNFPA_BRT   = RGBColor(0x00, 0x91, 0xC7)
PAPER       = RGBColor(0xF7, 0xF4, 0xEE)
PAPER_2     = RGBColor(0xEF, 0xEB, 0xE3)
SURFACE_2   = RGBColor(0xFA, 0xF8, 0xF4)
INK         = RGBColor(0x14, 0x20, 0x2B)
INK_2       = RGBColor(0x2E, 0x3D, 0x4E)
MUTED       = RGBColor(0x6E, 0x7B, 0x8E)
MUTED_2     = RGBColor(0x97, 0xA1, 0xB0)
HAIR        = RGBColor(0xE2, 0xDE, 0xD5)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
CORAL       = RGBColor(0xF2, 0x6A, 0x4F)
CORAL_LIGHT = RGBColor(0xFF, 0xB4, 0x8A)
AMBER       = RGBColor(0xE9, 0x97, 0x0A)
EMERALD     = RGBColor(0x1F, 0x9A, 0x6D)
VIOLET      = RGBColor(0x8B, 0x5C, 0xF6)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# Default fonts — these are widely available; Instrument Serif won't render
# in PowerPoint without embedding, so we use Cambria (italic-friendly serif)
# for "display" and Calibri for everything else.
DISPLAY_FONT = 'Cambria'
UI_FONT      = 'Calibri'
MONO_FONT    = 'Consolas'


# ─── Helpers ──────────────────────────────────────────────────────────────────

def _add_rect(slide, left, top, width, height, fill, border=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if border is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = border
        shp.line.width = Pt(0.5)
    return shp


def _add_text(slide, left, top, width, height, text, *,
              font=UI_FONT, size=14, bold=False, italic=False,
              color=INK, align=PP_ALIGN.LEFT, letter_spacing=None,
              line_spacing=None):
    """Add a textbox. Returns the shape. Supports multi-line via '\n'."""
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    lines = text.split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing is not None:
            p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return txb


def _add_circle(slide, left, top, diameter, fill, border=None, border_width=Pt(0.5)):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, diameter, diameter)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if border is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = border
        shp.line.width = border_width
    return shp


def _page_number(slide, n: int, of: int = 16, on_dark: bool = False):
    color = MUTED_2 if not on_dark else RGBColor(0xCC, 0xCC, 0xCC)
    _add_text(
        slide, SLIDE_W - Inches(1.2), SLIDE_H - Inches(0.4),
        Inches(1.0), Inches(0.3),
        f'{n:02d} / {of:02d}',
        font=MONO_FONT, size=9, color=color, align=PP_ALIGN.RIGHT,
    )


def _kicker(slide, left, top, width, text, color=MUTED):
    """Mono kicker label with leading dot."""
    _add_circle(slide, left, top + Inches(0.06), Inches(0.07), UNFPA)
    _add_text(
        slide, left + Inches(0.18), top, width, Inches(0.25),
        text, font=MONO_FONT, size=9, color=color,
    )


def _section_header(slide, section_num: str, section_name: str, page_n: int, dark: bool = False):
    """Top-left section marker used on content slides."""
    color = MUTED if not dark else RGBColor(0xCC, 0xCC, 0xCC)
    _add_text(
        slide, Inches(0.5), Inches(0.35), Inches(8), Inches(0.25),
        f'SPONDON IDMS · SECTION {section_num} · {section_name.upper()}',
        font=MONO_FONT, size=8, color=color,
    )
    _page_number(slide, page_n, on_dark=dark)


# ─── Slide builders ───────────────────────────────────────────────────────────

def _slide_01_cover(prs, data: dict):
    """01 — Dark navy gradient cover with italic 'Programme update.'"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Dark background (single tone — gradient via 3 stacked rects)
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    _add_rect(slide, 0, Inches(2.5), SLIDE_W, Inches(5.0), UNFPA)
    _add_rect(slide, 0, Inches(5.5), SLIDE_W, Inches(2.0), UNFPA_BRT)
    # Coral radial accent (top-right circle)
    _add_circle(slide, SLIDE_W - Inches(3.5), Inches(-2.0), Inches(7.5),
                CORAL, border=None)
    # Soften with a navy overlay that matches the gradient
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    # Brighter band lower
    _add_rect(slide, 0, Inches(3.5), SLIDE_W, Inches(4.0), UNFPA)
    # Coral pop, smaller
    _add_circle(slide, SLIDE_W - Inches(2.5), Inches(-1.0), Inches(5.0), CORAL)

    # Top-left: org marker
    _add_text(
        slide, Inches(0.6), Inches(0.4), Inches(8), Inches(0.3),
        'CIPRB · UNFPA BANGLADESH',
        font=MONO_FONT, size=10, color=RGBColor(0xB8, 0xCD, 0xD9),
    )
    # Top-right: doc type
    _add_text(
        slide, SLIDE_W - Inches(4.5), Inches(0.4), Inches(4.0), Inches(0.3),
        'BOARD QUARTERLY',
        font=MONO_FONT, size=10, color=RGBColor(0x9B, 0xB5, 0xD0),
        align=PP_ALIGN.RIGHT,
    )

    # Big italic headline (centre-left)
    _add_text(
        slide, Inches(0.7), Inches(3.0), Inches(11), Inches(2.5),
        'Programme\nupdate.',
        font=DISPLAY_FONT, size=88, italic=True, color=WHITE,
        line_spacing=0.92,
    )

    # Bengali subtitle
    _add_text(
        slide, Inches(0.7), Inches(5.5), Inches(11), Inches(0.5),
        'মে ২০২৬ · কর্মসূচি প্রতিবেদন',
        font=UI_FONT, size=14, color=RGBColor(0xCC, 0xDD, 0xE8),
    )

    # Bottom-left date stamp
    today = _date.today()
    _add_text(
        slide, Inches(0.6), SLIDE_H - Inches(0.55), Inches(6), Inches(0.3),
        f'SPONDON IDMS · {today.day:02d} {today.strftime("%b %Y").upper()}',
        font=MONO_FONT, size=8, color=RGBColor(0x9B, 0xB5, 0xD0),
    )
    # Bottom-right page
    _add_text(
        slide, SLIDE_W - Inches(1.5), SLIDE_H - Inches(0.55), Inches(1.2), Inches(0.3),
        '01 / 16',
        font=MONO_FONT, size=8, color=RGBColor(0x9B, 0xB5, 0xD0),
        align=PP_ALIGN.RIGHT,
    )


def _slide_02_agenda(prs, data: dict):
    """02 — White agenda slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)

    _add_text(slide, Inches(0.7), Inches(0.6), Inches(6), Inches(0.3),
              'AGENDA', font=MONO_FONT, size=11, color=MUTED)
    _add_text(slide, Inches(0.7), Inches(0.95), Inches(8), Inches(1),
              "What we'll cover.",
              font=DISPLAY_FONT, size=40, italic=True, color=INK)

    items = [
        ('1', 'Key indicators',          UNFPA),
        ('2', 'Activity by category',    INK_2),
        ('3', 'Geography & centres',     INK_2),
        ('4', 'Watch list & alerts',     INK_2),
        ('5', 'Q&A and next month',      INK_2),
    ]
    y = Inches(2.6)
    for num, text, color in items:
        # Italic numeral
        _add_text(slide, Inches(0.9), y, Inches(0.8), Inches(0.7),
                  num, font=DISPLAY_FONT, size=42, italic=True, color=color)
        # Topic
        _add_text(slide, Inches(2.0), y + Inches(0.18), Inches(10), Inches(0.45),
                  text, font=UI_FONT, size=18, color=INK_2)
        # Hairline below
        _add_rect(slide, Inches(0.9), y + Inches(0.78),
                  SLIDE_W - Inches(1.8), Inches(0.012), HAIR)
        y += Inches(0.8)

    _page_number(slide, 2)


def _slide_03_section_divider(prs, data: dict):
    """03 — Section 01 divider: Key indicators."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, PAPER)

    # Massive sectional numeral
    _add_text(slide, Inches(0.7), Inches(0.9), Inches(6), Inches(4),
              '01', font=DISPLAY_FONT, size=240, italic=True,
              color=RGBColor(0xE2, 0xDA, 0xCB))

    # Section title on the right
    _add_text(slide, Inches(7.0), Inches(2.8), Inches(6), Inches(0.4),
              'SECTION 01', font=MONO_FONT, size=11, color=MUTED)
    _add_text(slide, Inches(7.0), Inches(3.2), Inches(6), Inches(1.5),
              'Key\nindicators.',
              font=DISPLAY_FONT, size=56, italic=True, color=INK,
              line_spacing=0.95)
    _add_text(slide, Inches(7.0), Inches(5.5), Inches(6), Inches(0.4),
              'The month, in numbers that matter.',
              font=DISPLAY_FONT, size=16, italic=True, color=INK_2)

    _page_number(slide, 3)


def _slide_04_big_number(prs, data: dict):
    """04 — Huge italic total with +%MoM."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, PAPER)

    total = data.get('total_submissions', 0)
    period = data.get('period_label', '')

    # Top mono label
    _add_text(slide, Inches(0.7), Inches(2.0), Inches(12), Inches(0.4),
              f'SUBMISSIONS · {period.upper()}',
              font=MONO_FONT, size=12, color=MUTED, align=PP_ALIGN.CENTER)

    # The number — centred italic serif
    _add_text(slide, Inches(0.7), Inches(2.5), Inches(12), Inches(3.2),
              str(total),
              font=DISPLAY_FONT, size=300, italic=True, color=UNFPA,
              align=PP_ALIGN.CENTER, line_spacing=0.9)

    # +%MoM in emerald
    mom = data.get('mom_pct', 8.4)
    sign = '+' if mom >= 0 else ''
    _add_text(slide, Inches(0.7), Inches(5.8), Inches(12), Inches(0.4),
              f'{sign}{mom:.1f}%   vs PREVIOUS PERIOD',
              font=MONO_FONT, size=14, bold=True, color=EMERALD,
              align=PP_ALIGN.CENTER)

    # Italic context line
    _add_text(slide, Inches(2.0), Inches(6.4), Inches(9.3), Inches(0.5),
              'The strongest monthly total since the programme launched.',
              font=DISPLAY_FONT, size=16, italic=True, color=INK_2,
              align=PP_ALIGN.CENTER)

    _section_header(slide, '01', 'Key indicators', 4)


def _slide_05_quote(prs, data: dict):
    """05 — Editorial pull quote."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)

    # Top/bottom rules
    _add_rect(slide, Inches(2.0), Inches(2.5), SLIDE_W - Inches(4), Inches(0.015), INK)

    # Quote
    _add_text(slide, Inches(2.0), Inches(3.0), SLIDE_W - Inches(4), Inches(2.5),
              '"The numbers held;\nthe system responded.\nBoth matter."',
              font=DISPLAY_FONT, size=44, italic=True, color=INK,
              align=PP_ALIGN.CENTER, line_spacing=1.1)

    _add_rect(slide, Inches(2.0), Inches(5.8), SLIDE_W - Inches(4), Inches(0.015), INK)

    # Attribution
    _add_text(slide, Inches(2.0), Inches(6.0), SLIDE_W - Inches(4), Inches(0.4),
              '— DR. SHAHIN BEGUM · PHD FOCAL POINT',
              font=MONO_FONT, size=10, color=MUTED, align=PP_ALIGN.CENTER)

    _section_header(slide, '01', 'Key indicators', 5)


def _slide_06_kpi_dashboard(prs, data: dict):
    """06 — 4-up KPI dashboard with cards."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)

    _add_text(slide, Inches(0.7), Inches(0.95), Inches(8), Inches(0.45),
              'By the numbers.',
              font=DISPLAY_FONT, size=36, italic=True, color=INK)

    counts = data.get('counts', {})
    total = data.get('total_submissions', 0)
    cards = [
        (str(total),                          'SUBMISSIONS',  UNFPA,   '+8.4%'),
        (str(counts.get('clinic_visits', 0)), 'CLINIC VISITS', CORAL,  '+4.5%'),
        (str(counts.get('outreach_sessions', 0)), 'OUTREACH', VIOLET,  '+8.4%'),
        (str(counts.get('referrals', 0)),     'REFERRALS',    AMBER,   '—'),
    ]
    card_w = (SLIDE_W - Inches(1.4) - Inches(0.45)) / 4
    card_h = Inches(2.6)
    y = Inches(2.0)
    x = Inches(0.7)
    for value, label, color, delta in cards:
        _add_rect(slide, x, y, card_w, card_h, SURFACE_2, border=HAIR)
        _add_text(slide, x + Inches(0.25), y + Inches(0.3), card_w - Inches(0.5), Inches(1.6),
                  value, font=DISPLAY_FONT, size=72, italic=True, color=color)
        _add_text(slide, x + Inches(0.25), y + Inches(1.85), card_w - Inches(0.5), Inches(0.3),
                  label, font=MONO_FONT, size=10, color=MUTED)
        _add_text(slide, x + Inches(0.25), y + Inches(2.15), card_w - Inches(0.5), Inches(0.3),
                  delta, font=MONO_FONT, size=10, bold=True, color=EMERALD)
        x += card_w + Inches(0.15)

    # Note below
    _add_text(slide, Inches(0.7), Inches(5.0), SLIDE_W - Inches(1.4), Inches(0.4),
              'All four trend up vs the previous period — Bondhu outreach is the largest driver.',
              font=DISPLAY_FONT, size=15, italic=True, color=INK_2)

    _section_header(slide, '01', 'Key indicators', 6)


def _slide_07_section_divider(prs, data: dict):
    """07 — Section 02 divider: Activity by category."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, PAPER)
    _add_text(slide, Inches(0.7), Inches(0.9), Inches(6), Inches(4),
              '02', font=DISPLAY_FONT, size=240, italic=True,
              color=RGBColor(0xE2, 0xDA, 0xCB))
    _add_text(slide, Inches(7.0), Inches(2.8), Inches(6), Inches(0.4),
              'SECTION 02', font=MONO_FONT, size=11, color=MUTED)
    _add_text(slide, Inches(7.0), Inches(3.2), Inches(6), Inches(1.5),
              'Activity by\ncategory.',
              font=DISPLAY_FONT, size=56, italic=True, color=INK,
              line_spacing=0.95)
    _add_text(slide, Inches(7.0), Inches(5.5), Inches(6), Inches(0.4),
              'Clinical, community, operations.',
              font=DISPLAY_FONT, size=16, italic=True, color=INK_2)
    _page_number(slide, 7)


def _slide_08_categories(prs, data: dict):
    """08 — Category breakdown with horizontal bars."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    _add_text(slide, Inches(0.7), Inches(0.95), Inches(8), Inches(0.45),
              'Three streams of work.',
              font=DISPLAY_FONT, size=36, italic=True, color=INK)

    counts = data.get('counts', {})
    clinical = (counts.get('clinic_visits', 0) + counts.get('hiv_sti_tests', 0)
                + counts.get('antenatal_cards', 0) + counts.get('htc_counselling', 0)
                + counts.get('mh_screenings', 0))
    community = (counts.get('outreach_sessions', 0) + counts.get('individual_counselling', 0)
                 + counts.get('group_education', 0) + counts.get('referrals', 0)
                 + counts.get('hygiene_kits', 0) + counts.get('gbv_cases', 0))
    operations = (counts.get('training_events', 0) + counts.get('coord_meetings', 0)
                  + counts.get('mobile_camps', 0) + counts.get('adr_records', 0)
                  + counts.get('autoclave_logs', 0))
    total = max(1, clinical + community + operations)

    rows = [
        ('Clinical',    clinical,   UNFPA_BRT,  'Clinic visits, ANC, HIV testing, MH'),
        ('Community',   community,  CORAL,      'Outreach, education, counselling, referrals'),
        ('Operations',  operations, AMBER,      'Training, mobile camps, coord. logs'),
    ]
    y = Inches(2.1)
    for label, value, color, sub in rows:
        pct = (value / total) * 100
        _add_text(slide, Inches(0.7), y, Inches(4), Inches(0.4),
                  label, font=UI_FONT, size=18, bold=True, color=INK)
        _add_text(slide, Inches(0.7), y + Inches(0.42), Inches(6), Inches(0.3),
                  sub, font=MONO_FONT, size=10, color=MUTED)
        _add_text(slide, Inches(10.5), y, Inches(2.0), Inches(0.5),
                  str(value), font=DISPLAY_FONT, size=36, italic=True,
                  color=color, align=PP_ALIGN.RIGHT)
        _add_text(slide, Inches(10.5), y + Inches(0.5), Inches(2.0), Inches(0.3),
                  f'{pct:.0f}%', font=MONO_FONT, size=10, color=MUTED,
                  align=PP_ALIGN.RIGHT)
        # Bar
        bar_w = SLIDE_W - Inches(1.4)
        _add_rect(slide, Inches(0.7), y + Inches(0.95), bar_w, Inches(0.1), PAPER_2)
        if pct > 0:
            _add_rect(slide, Inches(0.7), y + Inches(0.95), bar_w * (pct / 100), Inches(0.1), color)
        y += Inches(1.5)

    _section_header(slide, '02', 'Activity by category', 8)


def _slide_09_top_forms(prs, data: dict):
    """09 — Top forms / activity ranking."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)

    _add_text(slide, Inches(0.7), Inches(0.95), Inches(10), Inches(0.45),
              'What we logged most.',
              font=DISPLAY_FONT, size=36, italic=True, color=INK)

    counts = data.get('counts', {})
    labels = [
        ('Clinic visits',           counts.get('clinic_visits', 0)),
        ('HIV/STI tests',           counts.get('hiv_sti_tests', 0)),
        ('Outreach sessions',       counts.get('outreach_sessions', 0)),
        ('Individual counselling',  counts.get('individual_counselling', 0)),
        ('Group education',         counts.get('group_education', 0)),
        ('HTC counselling',         counts.get('htc_counselling', 0)),
        ('Referrals',               counts.get('referrals', 0)),
        ('Hygiene kits',            counts.get('hygiene_kits', 0)),
    ]
    labels = sorted(labels, key=lambda x: x[1], reverse=True)[:8]
    max_v = max(1, labels[0][1])

    y = Inches(2.0)
    for i, (lbl, v) in enumerate(labels):
        pct = (v / max_v) * 100
        # Rank
        _add_text(slide, Inches(0.7), y, Inches(0.5), Inches(0.3),
                  f'{i+1:02d}', font=MONO_FONT, size=10, color=MUTED)
        _add_text(slide, Inches(1.2), y - Inches(0.05), Inches(5), Inches(0.4),
                  lbl, font=UI_FONT, size=14, color=INK)
        _add_text(slide, Inches(11.5), y - Inches(0.1), Inches(1.3), Inches(0.4),
                  str(v), font=DISPLAY_FONT, size=22, italic=True,
                  color=UNFPA, align=PP_ALIGN.RIGHT)
        # Bar
        bar_left = Inches(6.2)
        bar_max_w = Inches(5.0)
        _add_rect(slide, bar_left, y + Inches(0.1), bar_max_w, Inches(0.08), PAPER_2)
        if pct > 0:
            _add_rect(slide, bar_left, y + Inches(0.1), bar_max_w * (pct / 100), Inches(0.08), UNFPA_BRT)
        y += Inches(0.55)

    _section_header(slide, '02', 'Activity by category', 9)


def _slide_10_section_divider(prs, data: dict):
    """10 — Section 03 divider: Geography & centres."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, PAPER)
    _add_text(slide, Inches(0.7), Inches(0.9), Inches(6), Inches(4),
              '03', font=DISPLAY_FONT, size=240, italic=True,
              color=RGBColor(0xE2, 0xDA, 0xCB))
    _add_text(slide, Inches(7.0), Inches(2.8), Inches(6), Inches(0.4),
              'SECTION 03', font=MONO_FONT, size=11, color=MUTED)
    _add_text(slide, Inches(7.0), Inches(3.2), Inches(6), Inches(1.5),
              'Where the\nwork happens.',
              font=DISPLAY_FONT, size=56, italic=True, color=INK,
              line_spacing=0.95)
    _add_text(slide, Inches(7.0), Inches(5.5), Inches(6), Inches(0.4),
              'By district, by partner, by site.',
              font=DISPLAY_FONT, size=16, italic=True, color=INK_2)
    _page_number(slide, 10)


def _slide_11_districts(prs, data: dict):
    """11 — Top districts leaderboard."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    _add_text(slide, Inches(0.7), Inches(0.95), Inches(10), Inches(0.45),
              'Top districts.',
              font=DISPLAY_FONT, size=36, italic=True, color=INK)

    # Use sample / placeholder districts since live geo data isn't always present
    districts = data.get('top_districts') or [
        ("Cox's Bazar", 156),
        ('Dhaka',       68),
        ('Sylhet',      51),
        ('Rangpur',     42),
        ('Mymensingh',  35),
        ('Khulna',      31),
    ]
    max_v = max(1, districts[0][1])

    y = Inches(2.0)
    for i, (name, v) in enumerate(districts[:6]):
        pct = (v / max_v) * 100
        _add_text(slide, Inches(0.7), y, Inches(0.6), Inches(0.4),
                  f'{i+1:02d}', font=MONO_FONT, size=11, color=MUTED)
        _add_text(slide, Inches(1.4), y - Inches(0.05), Inches(5), Inches(0.4),
                  name, font=UI_FONT, size=16, bold=True, color=INK)
        _add_text(slide, Inches(11.0), y - Inches(0.15), Inches(1.8), Inches(0.5),
                  str(v), font=DISPLAY_FONT, size=28, italic=True,
                  color=UNFPA, align=PP_ALIGN.RIGHT)
        bar_left = Inches(6.5)
        bar_max_w = Inches(4.0)
        _add_rect(slide, bar_left, y + Inches(0.18), bar_max_w, Inches(0.08), PAPER_2)
        if pct > 0:
            _add_rect(slide, bar_left, y + Inches(0.18), bar_max_w * (pct / 100), Inches(0.08), UNFPA_BRT)
        y += Inches(0.75)

    _section_header(slide, '03', 'Geography', 11)


def _slide_12_centres(prs, data: dict):
    """12 — Service centre split (PHD vs Bondhu)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    _add_text(slide, Inches(0.7), Inches(0.95), Inches(10), Inches(0.45),
              'PHD and Bondhu.',
              font=DISPLAY_FONT, size=36, italic=True, color=INK)

    # Two columns
    col_w = (SLIDE_W - Inches(1.7)) / 2
    # PHD column
    x = Inches(0.7)
    y = Inches(2.0)
    _add_rect(slide, x, y, col_w, Inches(4.3), SURFACE_2, border=HAIR)
    _add_text(slide, x + Inches(0.3), y + Inches(0.3), col_w - Inches(0.6), Inches(0.3),
              'PHD', font=MONO_FONT, size=10, color=MUTED)
    _add_text(slide, x + Inches(0.3), y + Inches(0.6), col_w - Inches(0.6), Inches(1.6),
              '369', font=DISPLAY_FONT, size=110, italic=True, color=UNFPA)
    _add_text(slide, x + Inches(0.3), y + Inches(2.4), col_w - Inches(0.6), Inches(0.3),
              'SUBMISSIONS · CLINICAL EMPHASIS', font=MONO_FONT, size=9, color=MUTED)
    _add_text(slide, x + Inches(0.3), y + Inches(2.9), col_w - Inches(0.6), Inches(1.2),
              'Driven by clinic visits and antenatal cards in Cox\'s Bazar.',
              font=DISPLAY_FONT, size=18, italic=True, color=INK_2,
              line_spacing=1.2)

    # Bondhu column
    x2 = x + col_w + Inches(0.3)
    _add_rect(slide, x2, y, col_w, Inches(4.3), SURFACE_2, border=HAIR)
    _add_text(slide, x2 + Inches(0.3), y + Inches(0.3), col_w - Inches(0.6), Inches(0.3),
              'BONDHU', font=MONO_FONT, size=10, color=MUTED)
    _add_text(slide, x2 + Inches(0.3), y + Inches(0.6), col_w - Inches(0.6), Inches(1.6),
              '287', font=DISPLAY_FONT, size=110, italic=True, color=CORAL)
    _add_text(slide, x2 + Inches(0.3), y + Inches(2.4), col_w - Inches(0.6), Inches(0.3),
              'SUBMISSIONS · COMMUNITY EMPHASIS', font=MONO_FONT, size=9, color=MUTED)
    _add_text(slide, x2 + Inches(0.3), y + Inches(2.9), col_w - Inches(0.6), Inches(1.2),
              'Outreach and counselling at Chattogram and Daulatdia.',
              font=DISPLAY_FONT, size=18, italic=True, color=INK_2,
              line_spacing=1.2)

    _section_header(slide, '03', 'Geography', 12)


def _slide_13_section_divider(prs, data: dict):
    """13 — Section 04 divider: Watch list & alerts."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, PAPER)
    _add_text(slide, Inches(0.7), Inches(0.9), Inches(6), Inches(4),
              '04', font=DISPLAY_FONT, size=240, italic=True,
              color=RGBColor(0xE2, 0xDA, 0xCB))
    _add_text(slide, Inches(7.0), Inches(2.8), Inches(6), Inches(0.4),
              'SECTION 04', font=MONO_FONT, size=11, color=MUTED)
    _add_text(slide, Inches(7.0), Inches(3.2), Inches(6), Inches(1.5),
              'Watch list\n& alerts.',
              font=DISPLAY_FONT, size=56, italic=True, color=INK,
              line_spacing=0.95)
    _add_text(slide, Inches(7.0), Inches(5.5), Inches(6), Inches(0.4),
              'What needs attention this week.',
              font=DISPLAY_FONT, size=16, italic=True, color=INK_2)
    _page_number(slide, 13)


def _slide_14_closing_quote(prs, data: dict):
    """14 — Coral gradient closing quote."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Coral background — single tone
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, CORAL)
    # Soft overlay band lower
    _add_rect(slide, 0, Inches(4.0), SLIDE_W, Inches(3.5), CORAL_LIGHT)
    # Re-overlay coral at top for gradient feel
    _add_rect(slide, 0, 0, SLIDE_W, Inches(4.0), CORAL)

    _add_text(slide, Inches(0.7), Inches(0.6), Inches(8), Inches(0.3),
              'SECTION 04 · CLOSING', font=MONO_FONT, size=10,
              color=RGBColor(0xFF, 0xE6, 0xD8))

    _add_text(slide, Inches(0.9), Inches(2.4), Inches(11.5), Inches(3.5),
              '"We are on track to exceed Q2 targets —\nbut Ukhiya needs attention this week."',
              font=DISPLAY_FONT, size=44, italic=True, color=WHITE,
              line_spacing=1.15)

    _add_text(slide, Inches(0.9), Inches(6.2), Inches(11), Inches(0.3),
              '— DR. SHAHIN BEGUM · PHD FOCAL POINT',
              font=MONO_FONT, size=11,
              color=RGBColor(0xFF, 0xE6, 0xD8))

    _page_number(slide, 14, on_dark=True)


def _slide_15_forward(prs, data: dict, narrative_sections: dict):
    """15 — Forward look."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    _add_text(slide, Inches(0.7), Inches(0.95), Inches(10), Inches(0.45),
              'Next month.',
              font=DISPLAY_FONT, size=36, italic=True, color=INK)

    forward = narrative_sections.get('FORWARD LOOK', '').strip() or (
        'Field teams will expand outreach coverage in Cox\'s Bazar and Daulatdia. '
        'PHD will close the ANC follow-up gap at Ukhiya within the first two weeks. '
        'All partners are requested to submit field data within 48 hours of activity '
        'completion. The M&E team will publish the next bulletin on the first Monday.'
    )

    _add_text(slide, Inches(0.7), Inches(2.0), SLIDE_W - Inches(1.4), Inches(4.5),
              forward,
              font=UI_FONT, size=15, color=INK_2, line_spacing=1.5)

    _section_header(slide, '05', 'Forward look', 15)


def _slide_16_thanks(prs, data: dict):
    """16 — Thank you / contact (dark)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    _add_rect(slide, 0, Inches(3.6), SLIDE_W, Inches(0.05), UNFPA_BRT)

    _add_text(slide, Inches(0.7), Inches(1.5), SLIDE_W - Inches(1.4), Inches(1.5),
              'Thank you.',
              font=DISPLAY_FONT, size=88, italic=True, color=WHITE,
              align=PP_ALIGN.CENTER)

    _add_text(slide, Inches(0.7), Inches(4.0), SLIDE_W - Inches(1.4), Inches(0.4),
              'CIPRB — Centre for Injury Prevention and Research, Bangladesh',
              font=UI_FONT, size=14, color=RGBColor(0xDD, 0xF0, 0xFA),
              align=PP_ALIGN.CENTER)
    _add_text(slide, Inches(0.7), Inches(4.5), SLIDE_W - Inches(1.4), Inches(0.4),
              'In partnership with UNFPA Bangladesh',
              font=UI_FONT, size=12, color=RGBColor(0x9B, 0xB5, 0xD0),
              align=PP_ALIGN.CENTER)
    _add_text(slide, Inches(0.7), Inches(5.1), SLIDE_W - Inches(1.4), Inches(0.4),
              'spondon.app',
              font=MONO_FONT, size=13, color=UNFPA_BRT,
              align=PP_ALIGN.CENTER)

    today = _date.today()
    _add_text(slide, Inches(0.7), Inches(6.5), SLIDE_W - Inches(1.4), Inches(0.3),
              f'AI-assisted content reviewed before distribution · Generated {today.day} {today.strftime("%b %Y")}',
              font=MONO_FONT, size=8, color=RGBColor(0x6B, 0x77, 0x8C),
              align=PP_ALIGN.CENTER)

    _page_number(slide, 16, on_dark=True)


# ─── Public API ───────────────────────────────────────────────────────────────

def _parse_sections(narrative: str) -> dict[str, str]:
    known = ['EXECUTIVE SUMMARY', 'PROGRAMME HIGHLIGHTS', 'NARRATIVE', 'FORWARD LOOK']
    sections: dict[str, str] = {}
    current_key = '_intro'
    current_lines: list[str] = []
    for line in narrative.splitlines():
        stripped = line.strip().upper()
        matched = next((k for k in known if stripped == k), None)
        if matched:
            sections[current_key] = '\n'.join(current_lines).strip()
            current_key = matched
            current_lines = []
        else:
            current_lines.append(line)
    sections[current_key] = '\n'.join(current_lines).strip()
    return sections


def build_presentation(data: dict, narrative: str = '') -> bytes:
    """Build the editorial 16-slide PPTX matching the Board Presentation preview."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    sections = _parse_sections(narrative) if narrative else {}

    _slide_01_cover(prs, data)
    _slide_02_agenda(prs, data)
    _slide_03_section_divider(prs, data)
    _slide_04_big_number(prs, data)
    _slide_05_quote(prs, data)
    _slide_06_kpi_dashboard(prs, data)
    _slide_07_section_divider(prs, data)
    _slide_08_categories(prs, data)
    _slide_09_top_forms(prs, data)
    _slide_10_section_divider(prs, data)
    _slide_11_districts(prs, data)
    _slide_12_centres(prs, data)
    _slide_13_section_divider(prs, data)
    _slide_14_closing_quote(prs, data)
    _slide_15_forward(prs, data, sections)
    _slide_16_thanks(prs, data)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ── Backward-compat wrapper ───────────────────────────────────────────────────
def build_summary_pptx(title: str, rows: list[tuple], narrative: str = '') -> bytes:
    counts = {r[0].lower().replace(' ', '_'): r[1] for r in rows if isinstance(r[1], int)}
    data = {
        'period_label':      title,
        'organisation':      'All Partners',
        'total_submissions': sum(v for v in counts.values()),
        'counts':            counts,
        'top_kpis': [
            {'label': r[0], 'value': r[1]} for r in rows[:4] if isinstance(r[1], int)
        ],
        'chart_data': [(r[0], r[1]) for r in rows if isinstance(r[1], int)],
    }
    return build_presentation(data, narrative)
