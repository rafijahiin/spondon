"""
The board DECK — native, editable PowerPoint.

The old deck was seven full-bleed screenshots of the web template: nothing on
a slide could be edited, and it read as the report cut into pieces. This one
is built shape by shape with python-pptx — real text boxes, real tables, real
bars — so CIPRB/UNFPA staff can lift any slide into their own presentations
and fix a word without asking anyone.

Slide language (16:9, English per the 2026-08 decision):
  1  title           month, scope, brand
  2  the headline    hero number + what it means + the trend
  3  at a glance     KPI tiles as native shapes
  4+ one per block   label + horizontal bars
  n  geography       bars + coverage line          (when the scope has one)
  n  indicators      native table, target/achieved/%
  n  narrative       the AI paragraphs, attributed
  n  close           source line
"""
from __future__ import annotations

import io

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

from .brand import fmt

ORANGE = RGBColor(0xE8, 0x56, 0x2B)
ORANGE_DEEP = RGBColor(0xB2, 0x3A, 0x15)
INK = RGBColor(0x1A, 0x24, 0x2C)
MUTED = RGBColor(0x5C, 0x6C, 0x78)
FAINT = RGBColor(0x8A, 0x97, 0xA0)
SOFT = RGBColor(0xF6, 0xF4, 0xF1)
LINE = RGBColor(0xE3, 0xE0, 0xDA)
GREEN = RGBColor(0x2E, 0x7D, 0x54)
AMBER = RGBColor(0xC7, 0x7F, 0x1F)
RED = RGBColor(0xB6, 0x39, 0x2F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

W, H = Inches(13.333), Inches(7.5)
FONT = 'Atkinson Hyperlegible'


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _box(slide, x, y, w, h, fill=None, line=None, radius=False):
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        x, y, w, h)
    if radius:
        try:
            shape.adjustments[0] = 0.12
        except Exception:
            pass
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    return shape


def _text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, space_after=0):
    """runs: list of (text, size_pt, bold, colour) tuples → one paragraph each."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (txt, size, bold, colour) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if space_after:
            p.space_after = Pt(space_after)
        r = p.add_run()
        r.text = txt
        f = r.font
        f.name = FONT
        f.size = Pt(size)
        f.bold = bold
        f.color.rgb = colour
    return tb


def _brand(slide, page=None):
    from pptx.enum.shapes import MSO_SHAPE
    d = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.55), Inches(0.42),
                               Inches(0.16), Inches(0.16))
    d.fill.solid()
    d.fill.fore_color.rgb = ORANGE
    d.line.fill.background()
    d.shadow.inherit = False
    _text(slide, Inches(0.82), Inches(0.33), Inches(3), Inches(0.4),
          [('S I M P L E', 12, True, INK)])
    if page:
        _text(slide, W - Inches(1.2), H - Inches(0.5), Inches(0.8), Inches(0.35),
              [(str(page), 10, False, FAINT)], align=PP_ALIGN.RIGHT)


def _bars(slide, x, y, w, rows, row_h=Inches(0.52), label_w=Inches(3.6),
          val_w=Inches(1.0), accent=ORANGE):
    top = max((r['value'] for r in rows), default=0) or 1
    bar_w = w - label_w - val_w - Inches(0.3)
    for i, r in enumerate(rows):
        yy = y + i * row_h
        _text(slide, x, yy, label_w, row_h, [(r['en'], 13, False, INK)])
        track = _box(slide, x + label_w, yy + Inches(0.09), bar_w, Inches(0.2),
                     fill=SOFT, radius=True)
        frac = r['value'] / top
        if frac > 0:
            _box(slide, x + label_w, yy + Inches(0.09),
                 Emu(int(bar_w * frac)) if frac > 0.02 else Inches(0.12),
                 Inches(0.2), fill=accent, radius=True)
        _text(slide, x + label_w + bar_w + Inches(0.15), yy, val_w, row_h,
              [(fmt(r['value']), 14, True, INK)], align=PP_ALIGN.RIGHT)


def _title_slide(prs, c):
    s = _blank(prs)
    _box(s, 0, 0, W, H, fill=WHITE)
    _box(s, 0, H - Inches(0.35), W, Inches(0.35), fill=ORANGE)
    _brand(s)
    _text(s, Inches(0.55), Inches(2.2), Inches(11), Inches(0.5),
          [('MONTHLY PROGRAMME REVIEW', 13, True, ORANGE_DEEP)])
    _text(s, Inches(0.5), Inches(2.6), Inches(12), Inches(1.4),
          [(c['period_label'], 54, True, INK)])
    _text(s, Inches(0.55), Inches(4.0), Inches(12), Inches(0.5),
          [(f"{c['org_label']}   ·   CIPRB / UNFPA Bangladesh", 16, False, MUTED)])


def _hero_slide(prs, c):
    s = _blank(prs)
    _box(s, 0, 0, W, H, fill=WHITE)
    _brand(s, 2)
    hero = c['hero']
    _text(s, Inches(0.55), Inches(1.15), Inches(8), Inches(0.4),
          [('THE HEADLINE', 12, True, ORANGE_DEEP)])
    _text(s, Inches(0.45), Inches(1.7), Inches(7.5), Inches(2.2),
          [(fmt(hero['value']), 110, True, ORANGE)])
    _text(s, Inches(0.55), Inches(4.15), Inches(6.6), Inches(1.2),
          [(hero['en'], 20, True, INK), (hero['note'], 13, False, MUTED)],
          space_after=6)
    trend = c.get('trend') or []
    if trend and max(trend):
        bx, bw = Inches(8.2), Inches(4.4)
        top = max(trend)
        n = len(trend)
        col_w = int(bw / (n * 1.5))
        for i, v in enumerate(trend):
            hpx = int(Inches(2.2) * (v / top)) if v else Inches(0.03)
            _box(s, bx + int(i * bw / n), Inches(4.6) - hpx + Inches(0.0),
                 col_w, hpx, fill=ORANGE if i == n - 1 else LINE)
        _text(s, bx, Inches(4.75), bw, Inches(0.4),
              [('12-month trend of approved records', 11, False, FAINT)])
    if c.get('mom_pct') is not None:
        up = c['mom_pct'] >= 0
        _text(s, Inches(8.2), Inches(1.7), Inches(4), Inches(0.4),
              [(f"{'▲' if up else '▼'} {abs(c['mom_pct'])}% vs previous month",
                14, True, GREEN if up else ORANGE_DEEP)])


def _kpi_slide(prs, c):
    s = _blank(prs)
    _box(s, 0, 0, W, H, fill=WHITE)
    _brand(s, 3)
    _text(s, Inches(0.55), Inches(1.0), Inches(10), Inches(0.5),
          [('AT A GLANCE', 12, True, ORANGE_DEEP)])
    _text(s, Inches(0.55), Inches(1.35), Inches(10), Inches(0.6),
          [('The month in numbers', 26, True, INK)])
    kpis = c['kpis'][:6]
    cols = 3
    cw, ch = Inches(3.95), Inches(1.9)
    gx, gy = Inches(0.55), Inches(2.3)
    for i, k in enumerate(kpis):
        x = gx + (i % cols) * (cw + Inches(0.25))
        y = gy + (i // cols) * (ch + Inches(0.25))
        _box(s, x, y, cw, ch, fill=SOFT, radius=True)
        _text(s, x + Inches(0.25), y + Inches(0.22), cw - Inches(0.5), Inches(0.9),
              [(fmt(k['value']), 36, True, INK)])
        _text(s, x + Inches(0.25), y + Inches(1.15), cw - Inches(0.5), Inches(0.65),
              [(k['en'], 12.5, False, MUTED)])


def _block_slide(prs, c, block, page):
    s = _blank(prs)
    _box(s, 0, 0, W, H, fill=WHITE)
    _brand(s, page)
    _text(s, Inches(0.55), Inches(1.0), Inches(10), Inches(0.5),
          [('PROGRAMME DETAIL', 12, True, ORANGE_DEEP)])
    _text(s, Inches(0.55), Inches(1.35), Inches(11), Inches(0.6),
          [(block['en'], 26, True, INK)])
    _bars(s, Inches(0.7), Inches(2.35), Inches(11.9), block['rows'][:8])


def _geo_slide(prs, c, page):
    s = _blank(prs)
    _box(s, 0, 0, W, H, fill=WHITE)
    _brand(s, page)
    _text(s, Inches(0.55), Inches(1.0), Inches(10), Inches(0.5),
          [('WHERE', 12, True, ORANGE_DEEP)])
    _text(s, Inches(0.55), Inches(1.35), Inches(11), Inches(0.6),
          [(c['geo']['en'], 26, True, INK)])
    rows = [{'en': n, 'value': v} for n, v in c['geo']['rows']]
    _bars(s, Inches(0.7), Inches(2.35), Inches(11.9), rows[:7])
    if c['geo'].get('coverage'):
        _text(s, Inches(0.7), H - Inches(1.0), Inches(10), Inches(0.4),
              [(c['geo']['coverage'], 13, True, MUTED)])


def _indicator_slide(prs, c, page):
    s = _blank(prs)
    _box(s, 0, 0, W, H, fill=WHITE)
    _brand(s, page)
    _text(s, Inches(0.55), Inches(0.95), Inches(10), Inches(0.5),
          [('AGAINST THE M&E FRAMEWORK', 12, True, ORANGE_DEEP)])
    _text(s, Inches(0.55), Inches(1.3), Inches(11), Inches(0.6),
          [('Indicator progress', 26, True, INK)])
    rows = c['indicators'][:8]
    tbl = prs.slides[-1].shapes if False else None
    shape = s.shapes.add_table(len(rows) + 1, 4, Inches(0.55), Inches(2.15),
                               Inches(12.2), Inches(0.4) * (len(rows) + 1))
    t = shape.table
    for j, htxt in enumerate(('Indicator', 'Target', 'Achieved', 'Progress')):
        cell = t.cell(0, j)
        cell.text = htxt
        p = cell.text_frame.paragraphs[0]
        p.runs[0].font.size = Pt(12)
        p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = WHITE
        p.runs[0].font.name = FONT
        cell.fill.solid()
        cell.fill.fore_color.rgb = INK
    t.columns[0].width = Inches(7.0)
    for j in (1, 2, 3):
        t.columns[j].width = Inches(1.73)
    for i, r in enumerate(rows, start=1):
        vals = (f"{r['code']}  ·  {r['label'][:64]}",
                fmt(r['target']) if r['target'] is not None else '—',
                fmt(r['achieved']),
                f"{r['pct']:.0f}%" if r['pct'] is not None else 'no target')
        for j, v in enumerate(vals):
            cell = t.cell(i, j)
            cell.text = str(v)
            p = cell.text_frame.paragraphs[0]
            f = p.runs[0].font
            f.size = Pt(11.5)
            f.name = FONT
            f.bold = (j == 2)
            if j == 3:
                pcv = r['pct']
                f.color.rgb = (FAINT if pcv is None else
                               GREEN if pcv >= 75 else AMBER if pcv >= 40 else RED)
                f.bold = True
            else:
                f.color.rgb = INK
            if j > 0:
                p.alignment = PP_ALIGN.RIGHT
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if i % 2 else SOFT


def _narrative_slide(prs, paras, meta_line, page):
    s = _blank(prs)
    _box(s, 0, 0, W, H, fill=WHITE)
    _brand(s, page)
    _text(s, Inches(0.55), Inches(1.0), Inches(10), Inches(0.5),
          [('NARRATIVE', 12, True, ORANGE_DEEP)])
    runs = [(p, 16, False, INK) for p in paras[:3]]
    _text(s, Inches(0.55), Inches(1.7), Inches(12.2), Inches(4.4), runs,
          space_after=14)
    _text(s, Inches(0.55), H - Inches(1.0), Inches(12), Inches(0.4),
          [(f'✦ {meta_line}', 11, False, MUTED)])


def _close_slide(prs, c):
    s = _blank(prs)
    _box(s, 0, 0, W, H, fill=INK)
    _box(s, 0, H - Inches(0.35), W, Inches(0.35), fill=ORANGE)
    _text(s, Inches(0.55), Inches(2.9), Inches(12.2), Inches(1.2),
          [('Every figure in this deck is an approved submission', 28, True, WHITE),
           ('in SIMPLE — the same numbers as the live dashboard.', 28, True, WHITE)])
    _text(s, Inches(0.55), Inches(4.6), Inches(12), Inches(0.5),
          [('simpledashboard.pro  ·  CIPRB / UNFPA Bangladesh', 14, False,
            RGBColor(0x9D, 0xB4, 0xC0))])


def build_native_deck(c: dict, paras: list[str], meta_line: str) -> bytes:
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    _title_slide(prs, c)
    _hero_slide(prs, c)
    page = 3
    if c['kpis']:
        _kpi_slide(prs, c)
        page += 1
    for b in c['blocks']:
        if any(r['value'] for r in b['rows']):
            _block_slide(prs, c, b, page)
            page += 1
    if c.get('partners'):
        for p in c['partners']:
            d = p['data']
            rows = [k for k in d['kpis'][:5]]
            _block_slide(prs, c, {'en': d['org_label'], 'rows': rows}, page)
            page += 1
    if c['geo']['rows']:
        _geo_slide(prs, c, page)
        page += 1
    if c['indicators']:
        _indicator_slide(prs, c, page)
        page += 1
    if paras:
        _narrative_slide(prs, paras, meta_line, page)
        page += 1
    _close_slide(prs, c)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
