"""
HTML-first report rendering. The branded deliverables (infographic, report,
PPTX slides, web report) are authored as HTML/CSS templates and rendered to
PNG/PDF with headless Chromium (Playwright) — pixel-identical to the browser,
so screen, paper, deck and the live link all share one design language.

  render_infographic_png(data) -> PNG bytes   (the Monthly Programme Pulse)

`data` is the dict from reports.generators.data.collect_programme_data(); the
same dict feeds every format, so "pull fresh data each month" is a direct bind.
"""
from __future__ import annotations

from datetime import date, timedelta

from .data import LABEL_MAP

ORG_LINE = 'CIPRB · UNFPA Bangladesh'
PARTNER_COLORS = {'PHD': '#F96000', 'Bandhu': '#FF9A52', 'CIPRB': '#C44E00'}


# ── Playwright render primitives ────────────────────────────────────────────
def _launch(p):
    # --no-sandbox is required in the Railway/Docker container (no user ns).
    return p.chromium.launch(args=['--no-sandbox', '--disable-dev-shm-usage'])


def html_to_png(html: str, selector: str = '.sheet', scale: int = 2,
                wait_ms: int = 1400) -> bytes:
    """Render an HTML string and screenshot one element to PNG bytes."""
    from playwright.sync_api import sync_playwright
    with sync_playwright() as p:
        b = _launch(p)
        try:
            pg = b.new_page(device_scale_factor=scale)
            pg.set_content(html, wait_until='load')
            try:
                pg.wait_for_load_state('networkidle', timeout=8000)
            except Exception:
                pass
            pg.wait_for_timeout(wait_ms)           # web font + load animation
            return pg.locator(selector).screenshot()
        finally:
            b.close()


def html_to_pdf(html: str, wait_ms: int = 1400, **pdf_kw) -> bytes:
    """Render an HTML string to a print-quality PDF (bytes)."""
    from playwright.sync_api import sync_playwright
    with sync_playwright() as p:
        b = _launch(p)
        try:
            pg = b.new_page()
            pg.set_content(html, wait_until='load')
            try:
                pg.wait_for_load_state('networkidle', timeout=8000)
            except Exception:
                pass
            pg.wait_for_timeout(wait_ms)
            opts = dict(format='A4', print_background=True,
                        margin={'top': '0', 'bottom': '0', 'left': '0', 'right': '0'})
            opts.update(pdf_kw)
            return pg.pdf(**opts)
        finally:
            b.close()


def _sparkline(trend, w: int = 170, h: int = 48, pad: int = 4):
    """12-point polyline for the hero sparkline → (points, last_x, last_y)."""
    vals = list(trend or [])
    if len(vals) < 2:
        return '', 0, 0
    lo, hi = min(vals), max(vals)
    rng = (hi - lo) or 1
    n = len(vals)
    pts, lx, ly = [], 0, 0
    for i, v in enumerate(vals):
        x = pad + (w - 2 * pad) * i / (n - 1)
        y = h - pad - (h - 2 * pad) * (v - lo) / rng
        pts.append(f'{x:.0f},{y:.0f}')
        lx, ly = round(x), round(y)
    return ' '.join(pts), lx, ly


# ── Infographic: data dict → template context ───────────────────────────────
def infographic_context(data: dict, *, ai_summary: str = '', is_sample: bool = False,
                        districts_total: int = 19, generated_date: str | None = None) -> dict:
    from django.utils import timezone

    total = data.get('total_submissions', 0)
    org = data.get('organisation') or 'All Partners'
    is_overall = (org == 'All Partners')

    pe = data.get('period_end')
    period_title = pe.strftime('%B %Y') if pe else ''
    prev_month = ''
    if data.get('period_start'):
        prev_month = (data['period_start'] - timedelta(days=1)).strftime('%b')
    mom = data.get('mom_pct', 0) or 0

    # Partner split only makes sense on the overall sheet.
    breakdown = None
    if is_overall:
        bp = data.get('by_partner', {}) or {}
        mx = max(bp.values()) if bp else 1
        mx = mx or 1
        breakdown = []
        for name in ('PHD', 'Bandhu', 'CIPRB'):
            v = bp.get(name, 0)
            share = round(v / total * 100) if total else 0
            breakdown.append({'name': name, 'value_fmt': f'{v:,}',
                              'pct': round(v / mx * 100), 'color': PARTNER_COLORS[name],
                              'sub': f'{share}% of activity'})

    counts = data.get('counts', {}) or {}
    top = sorted(counts.items(), key=lambda kv: kv[1], reverse=True)[:6]
    tiles = [{'value_fmt': f'{v:,}', 'label': LABEL_MAP.get(k, k)} for k, v in top]

    td = data.get('top_districts', []) or []
    mxd = max((n for _, n in td), default=1) or 1
    districts = [{'name': n, 'value_fmt': f'{v:,}', 'pct': round(v / mxd * 100)} for n, v in td]

    sp, lx, ly = _sparkline(data.get('monthly_trend', []))
    focus = ('<b>All Partners</b> · PHD · Bandhu · CIPRB &nbsp;|&nbsp; Reproductive &amp; Child Health'
             if is_overall else f'<b>{org}</b> &nbsp;|&nbsp; Reproductive &amp; Child Health')

    return {
        'is_sample': is_sample, 'org': org, 'org_line': ORG_LINE, 'kicker': 'Monthly Programme Pulse',
        'period_title': period_title, 'focus_line': focus, 'ai_summary': ai_summary,
        'total_fmt': f'{total:,}',
        'districts_covered': data.get('districts_active') or len(td), 'districts_total': districts_total,
        'mom_abs': abs(mom), 'mom_dir': '▲' if mom >= 0 else '▼',
        'mom_class': 'pos' if mom >= 0 else 'neg', 'prev_month': prev_month,
        'spark_points': sp, 'spark_last_x': lx, 'spark_last_y': ly,
        'breakdown': breakdown, 'breakdown_title': 'By partner', 'breakdown_note': 'share of activity',
        'tiles': tiles, 'districts': districts,
        'fistula_fmt': f"{data.get('fistula_cases', 0):,}",
        'mpdsr_fmt': f"{data.get('mpdsr_cases', 0):,}",
        'workers_fmt': f"{data.get('active_workers', 0):,}",
        'foot_note': ('Narrative drafted by AI, figures from approved submissions only'
                      if ai_summary else 'Figures from approved submissions only'),
        'generated_date': generated_date or timezone.now().strftime('%d %b %Y'),
    }


def render_infographic_png(data: dict, **kw) -> bytes:
    from django.template.loader import render_to_string
    html = render_to_string('reports/infographic.html', infographic_context(data, **kw))
    return html_to_png(html, '.sheet')


# ── Full programme report: data dict → multi-page PDF ───────────────────────
def _trend_bars(vals, period_end):
    """12-month vertical bars with month labels."""
    vals = list(vals or [])
    if not vals:
        return []
    mx = max(vals) or 1
    labels = []
    if period_end:
        idx0 = period_end.year * 12 + (period_end.month - 1)
        for i in range(len(vals) - 1, -1, -1):
            ti = idx0 - i
            labels.append(date(ti // 12, ti % 12 + 1, 1).strftime('%b'))
    else:
        labels = [''] * len(vals)
    return [{'height': max(4, round(v / mx * 100)), 'label': lab}
            for v, lab in zip(vals, labels)]


def report_context(data: dict, *, narrative=None, closing_line=None, **kw) -> dict:
    ctx = infographic_context(data, **kw)          # reuse the shared brand fields
    counts = data.get('counts', {}) or {}
    total = data.get('total_submissions', 0)
    ctx['narrative'] = narrative or []
    ctx['kpis'] = [
        {'value_fmt': f'{total:,}', 'label': 'Total activities'},
        {'value_fmt': f"{counts.get('registrations', 0):,}", 'label': 'Clients registered'},
        {'value_fmt': f"{counts.get('individual_counselling', 0):,}", 'label': 'Counselling sessions'},
        {'value_fmt': f"{counts.get('referrals', 0):,}", 'label': 'Referrals'},
        {'value_fmt': f"{data.get('active_workers', 0):,}", 'label': 'Active field workers'},
        {'value_fmt': f"{data.get('pending', 0):,}", 'label': 'Pending review'},
    ]
    rows = sorted(((LABEL_MAP.get(k, k), v) for k, v in counts.items() if v > 0),
                  key=lambda x: x[1], reverse=True)
    mx = rows[0][1] if rows else 1
    ctx['service_rows'] = [{'label': l, 'value_fmt': f'{v:,}', 'pct': round(v / mx * 100)}
                           for l, v in rows]
    ctx['trend'] = _trend_bars(data.get('monthly_trend', []), data.get('period_end'))
    ctx['closing_line'] = closing_line or ('Strong momentum into next month — sustaining '
                                           'reach and closing the open action items.')
    return ctx


def render_report_pdf(data: dict, *, narrative=None, **kw) -> bytes:
    from django.template.loader import render_to_string
    html = render_to_string('reports/report.html', report_context(data, narrative=narrative, **kw))
    return html_to_pdf(html)


# ── PowerPoint: render each slide to PNG, embed full-bleed into a .pptx ──────
def render_slides_pngs(html: str, selector: str = '.slide', scale: int = 2,
                       wait_ms: int = 1400) -> list[bytes]:
    from playwright.sync_api import sync_playwright
    pngs: list[bytes] = []
    with sync_playwright() as p:
        b = _launch(p)
        try:
            pg = b.new_page(device_scale_factor=scale)
            pg.set_content(html, wait_until='load')
            try:
                pg.wait_for_load_state('networkidle', timeout=8000)
            except Exception:
                pass
            pg.wait_for_timeout(wait_ms)
            for el in pg.locator(selector).all():
                pngs.append(el.screenshot())
        finally:
            b.close()
    return pngs


def render_pptx(data: dict, *, narrative=None, **kw) -> bytes:
    import io
    from django.template.loader import render_to_string
    from pptx import Presentation
    from pptx.util import Inches
    html = render_to_string('reports/slides.html', report_context(data, narrative=narrative, **kw))
    pngs = render_slides_pngs(html, '.slide')
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for png in pngs:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(io.BytesIO(png), 0, 0,
                                 width=prs.slide_width, height=prs.slide_height)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ── Interactive web report (served live, shareable link) ────────────────────
def web_report_html(data: dict, *, narrative=None, **kw) -> str:
    """The animated monthly web report as an HTML string (served by a view)."""
    from django.template.loader import render_to_string
    return render_to_string('reports/web_report.html',
                            report_context(data, narrative=narrative, **kw))


# ── Demo dataset (collect_programme_data shape) for local proofs ─────────────
def demo_data() -> dict:
    return {
        'period_start': date(2026, 6, 1), 'period_end': date(2026, 6, 30),
        'period_label': '1–30 Jun 2026', 'organisation': 'All Partners',
        'total_submissions': 4820,
        'counts': {
            'registrations': 3512, 'individual_counselling': 2067, 'hiv_sti_tests': 1284,
            'antenatal_cards': 540, 'referrals': 311, 'gbv_cases': 96,
            'clinic_visits': 88, 'outreach_sessions': 74, 'mh_screenings': 41,
        },
        'fistula_cases': 27, 'mpdsr_cases': 41, 'active_workers': 38, 'pending': 19,
        'mom_pct': 18.0,
        'monthly_trend': [2900, 3100, 3050, 3360, 3520, 3700, 3980, 4120, 4300, 4490, 4610, 4820],
        'top_districts': [('Sunamganj', 812), ('Dhaka', 642), ('Habiganj', 518),
                          ('Noakhali', 423), ('Bhola', 357), ('Chandpur', 301)],
        'by_partner': {'PHD': 2140, 'Bandhu': 1760, 'CIPRB': 920},
    }
